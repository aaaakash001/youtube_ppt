import argparse
import cv2
from pptx import Presentation
from pptx.util import Inches
import numpy as np
from PIL import Image
from pytube import YouTube
from .utils import download_youtube_video, calculate_frame_difference, ocr_on_image
import os
from tqdm import tqdm

def main(video_path, output_pptx,percent=10.0):
    # Create 'ppt' directory if it doesn't exist
    ppt_dir = 'ppt'
    if not os.path.exists(ppt_dir):
        os.makedirs(ppt_dir)

    # Open the video file
    cap = cv2.VideoCapture(video_path)
    
    # Get total number of frames
    total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
    
    # Initialize variables
    ret, first_frame = cap.read()
    prev_frame = first_frame
    
    # Create a PowerPoint presentation with 16:9 aspect ratio
    presentation = Presentation()
    presentation.slide_width = Inches(16)
    presentation.slide_height = Inches(9)
    
    # Create progress bar
    progress_bar = tqdm(total=total_frames, desc="Processing Frames", unit="frame")
    
    count = 0
    # Read video frames and generate slides
    while True:
        count += 1
        ret, frame = cap.read()
        if not ret:
            break  # Break if no more frames are available
        
        # Calculate the percentage difference between frames
        diff_percentage = calculate_frame_difference(prev_frame, frame)
        
        if diff_percentage > float(percent):
            # Save the current frame as an image (PNG)
            cv2.imwrite(os.path.join(ppt_dir, 'temp_frame.png'), frame)
            img_path = os.path.join(ppt_dir, 'temp_frame.png')
            
            # Perform OCR on the image to extract text
            text = ocr_on_image(img_path)
            
            # Add a new slide to the presentation
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])  # Index 5 corresponds to a blank slide
            
            # Add the extracted text to the slide
            title = slide.shapes.title
            title.text = text
            
            # Add the current frame image to the slide
            slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=Inches(16), height=Inches(10))
            
        # Update the previous frame
        prev_frame = frame
        
        # Update progress bar
        progress_bar.update(1)
    
    # Close progress bar
    progress_bar.close()
    
    # Save the PowerPoint presentation inside 'ppt' directory
    presentation.save(os.path.join(ppt_dir, output_pptx))
    
    # Release the video capture object
    cap.release()
    
    # Delete temporary frame image
    os.remove(os.path.join(ppt_dir, 'temp_frame.png'))
    
    # Remove the temporary video file
    # os.remove(video_path)

if __name__ == "__main__":
    import sys
    args = sys.argv[1:]
    if len(args) !=3 :
        print("Usage: youtube-ocr <youtube_link> <output_pptx> similarity_percent")
        sys.exit(1)
    youtube_link, output_pptx,percent = args
    print("YouTube Link:", youtube_link)
    print("Output PowerPoint file:", output_pptx)
    download_youtube_video(youtube_link, 'temp_video')
    # Get the name of the downloaded file
    video_files = os.listdir('temp_video')
    if len(video_files) != 1:
        print("Error: Multiple or no video files found in 'temp_video' directory.")
        sys.exit(1)
    
    video_file_name = video_files[0]
    main(os.path.join('temp_video', video_file_name), output_pptx,percent)
