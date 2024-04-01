import argparse
import cv2
from pptx import Presentation
from pptx.util import Inches
import numpy as np
from PIL import Image
from pytube import YouTube
from .utils import download_youtube_video, calculate_frame_difference, ocr_on_image
import os

def main(video_path, output_pptx):
    # Open the video file
    cap = cv2.VideoCapture(video_path)
    
    # Initialize variables
    ret, first_frame = cap.read()
    prev_frame = first_frame
    
    # Create a PowerPoint presentation with 16:9 aspect ratio
    presentation = Presentation()
    presentation.slide_width = Inches(16)
    presentation.slide_height = Inches(9)
    
    count = 0
    # Read video frames and generate slides
    while True:
        count = count + 1
        ret, frame = cap.read()
        if not ret:
            break  # Break if no more frames are available
        
        # Calculate the percentage difference between frames
        diff_percentage = calculate_frame_difference(prev_frame, frame)
        
        if diff_percentage > 5.0:
            # Save the current frame as an image (PNG)
            cv2.imwrite('temp_frame.png', frame)
            img_path = 'temp_frame.png'
            
            # Perform OCR on the image to extract text
            text = ocr_on_image(img_path)
            
            # Add a new slide to the presentation
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])  # Index 5 corresponds to a blank slide
            
            # Add the extracted text to the slide
            title = slide.shapes.title
            title.text = text
            
            # Add the current frame image to the slide
            slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=Inches(16), height=Inches(10))
            
        # Update the previous frame
        prev_frame = frame
    
    # Save the PowerPoint presentation
    presentation.save(output_pptx)
    
    # Release the video capture object
    cap.release()
    
    # Delete temporary frame image
    os.remove('temp_frame.png')

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Download and process a YouTube video with OCR")
    parser.add_argument("youtube_link", help="YouTube video link")
    parser.add_argument("output_pptx", help="Output PowerPoint file")
    args = parser.parse_args()
    download_youtube_video(args.youtube_link, 'temp_video.mp4')
    main('temp_video.mp4', args.output_pptx)
